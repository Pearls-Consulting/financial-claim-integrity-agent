"""Build the "SDB Invoicing Solution — progress snapshot" deck from screenshots.

Screenshots are matched to slides by the leading number in the file name
(`1_home.png`, `5_boq_evidence.png`, ...). Two files with the same number go
side by side. Usage (from backend/):

    .venv/Scripts/python scripts/build_screenshots_deck.py [shots_dir]

shots_dir defaults to supporting_docs/presentation/shots/.
Output: supporting_docs/presentation/SDB-Invoicing-Solution-Progress.pptx
Slides with no matching screenshot get a grey placeholder naming the file
expected. Re-run after adding files; the deck is regenerated from scratch.
"""

from __future__ import annotations

import re
import sys
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
PRES_DIR = ROOT / "supporting_docs" / "presentation"
SHOTS = Path(sys.argv[1]) if len(sys.argv) > 1 else PRES_DIR / "shots"
OUT = PRES_DIR / "SDB-Invoicing-Solution-Progress.pptx"

NAVY = RGBColor(0x1B, 0x2A, 0x4A)
GREY = RGBColor(0x6B, 0x72, 0x80)
LIGHT = RGBColor(0xF3, 0xF4, 0xF6)
BORDER = RGBColor(0xD1, 0xD5, 0xDB)

NO_SHOT = "no screenshot"

# (slide no, title, caption lines, expected screenshot file stem or NO_SHOT)
SLIDES: list[tuple[int, str, list[str], str]] = [
    (0, "SDB Invoicing Solution — Claim Integrity Agent",
     [f"Progress snapshot · {date.today():%d %B %Y}"],
     NO_SHOT),
    (1, "Claims intake queue — استلام المطالبات",
     ["Vendor claims pulled from Dynamics 365, each showing where it is in the review and the agent's recommendation.",
      "• Fully bilingual — Arabic RTL / English with one toggle",
      "• Dynamics 365 integration — claims, PO and contract values come from the ERP"],
     "1_home"),
    (2, "Step 1 — Tax invoice intake",
     ["The reviewer uploads the vendor's tax invoice. The agent reads it and pre-fills the claim — invoice number, date, amounts and VAT — so nothing is re-typed.",
      "• Guided 6-step review, one step per gate of SDB's procedure (SP-01-04-05-02)",
      "• The same claim can be imported directly from Dynamics 365 instead"],
     "2_invoice"),
    (3, "Step 1 — ZATCA QR verification",
     ["The agent decodes the QR code printed on the tax invoice and compares every field it carries — seller, VAT number, timestamp, totals — against the invoice face. Here all five match and the Phase-2 digital signature verifies.",
      "• Each check cites its source in SDB's procedure or the regulation",
      "• Evidence is one click away: every compared value links to its location in the PDF"],
     "3_invoice_qr"),
    (4, "Step 2 — Invoice lines vs the contract's Bill of Quantities",
     ["The agent extracts every line from the invoice and matches it to the contracted BoQ — item code, unit, quantity and unit price. 7 of 12 contracted items are billed, every price matches the schedule, and the 5 unbilled items are flagged as normal for a periodic claim.",
      "• Also checks the payment sequence and the claim type against the payment history (a known D365 rejection reason at SDB)"],
     "4_boq"),
    (5, "Evidence behind every extraction",
     ["Every extracted value links back to its exact location in the source document — one click opens the invoice with the figure highlighted. Reviewers verify; they don't take the agent's word for it."],
     "5_boq_evidence"),
    (6, "Step 3 — Acceptance & three-way match",
     ["The agent confirms the right acceptance document exists for the contract type — here the completion certificate COC-000000355 — and that its certified amount equals the claim total. Where an ERP product receipt exists, quantities are reconciled line by line across contract, receipt and invoice."],
     "6_threeway"),
    (7, "Step 4 — Final check: completion certificate, dates and penalties",
     ["The completion certificate is cross-checked against the vendor's penalty register and the contract dates: delivery 165 days before the contract end, and the certificate's 'no delay / no stoppage' declarations agree with the penalty record.",
      "• Catches the case SDB reported — a certificate declaring no delay for a vendor who was penalised for one"],
     "7_penalties"),
    (8, "Step 5 — Vendor file: legal & compliance documents",
     ["The agent identifies each document in the vendor file — award letter, work commencement, commercial registration, GOSI and Zakat certificates — reads its reference number and validity date, and confirms the pack required by procedure step 6 is complete.",
      "• Completeness is based on what the agent actually read, not on a checkbox list",
      "• Expired or missing certificates are flagged before the claim reaches finance"],
     "8_prefinance"),
    (9, "Step 6 — Recommendation and hand-off",
     ["The agent issues its recommendation with a written rationale citing the procedure steps and regulations behind it, and a one-line verdict per gate. The reviewer keeps the decision.",
      "• Export bundles every document that took part in the matching — tax invoice, contract/BoQ, completion certificate — named by claim number, ready for finance and audit"],
     "9_summary_and_export"),
    (10, "Next steps",
     ["• Connect to SDB's Dynamics 365 tenant (claims, attachments, product receipts)",
      "• Contract reading for long (300-page) contracts — direct vision-model read under test",
      "• Cross-check CR and VAT numbers on the vendor certificates against the invoice and ERP vendor record",
      "• Payment recommendation: retention / penalties / net payable, with the vendor notification",
      "• Pilot with SDB's vendor-management team on live claims"],
     NO_SHOT),
]


def find_shots() -> dict[int, list[Path]]:
    found: dict[int, list[Path]] = {}
    if SHOTS.exists():
        for p in sorted(SHOTS.iterdir()):
            m = re.match(r"(\d+)", p.stem)
            if m and p.suffix.lower() in (".png", ".jpg", ".jpeg"):
                found.setdefault(int(m.group(1)), []).append(p)
    return found


def add_text(slide, left, top, width, height, lines, size, bold=False, color=NAVY, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        para.space_after = Pt(4)
        run = para.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def fit_picture(slide, path: Path, left, top, max_w, max_h):
    pic = slide.shapes.add_picture(str(path), left, top)
    ratio = min(max_w / pic.width, max_h / pic.height)
    pic.width, pic.height = int(pic.width * ratio), int(pic.height * ratio)
    pic.left = int(left + (max_w - pic.width) / 2)
    pic.top = int(top + (max_h - pic.height) / 2)
    pic.line.color.rgb = BORDER
    pic.line.width = Pt(0.75)


def placeholder(slide, left, top, w, h, hint: str):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = BORDER
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = f"SCREENSHOT TO ADD\n\n{hint}.png"
    r.font.size = Pt(14)
    r.font.color.rgb = GREY


def build() -> None:
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    W = prs.slide_width
    blank = prs.slide_layouts[6]
    shots = find_shots()
    missing: list[str] = []

    for no, title, caption, hint in SLIDES:
        s = prs.slides.add_slide(blank)
        if no == 0:
            add_text(s, Inches(0.8), Inches(2.4), W - Inches(1.6), Inches(1.4), [title], 36, True, align=PP_ALIGN.CENTER)
            add_text(s, Inches(0.8), Inches(3.9), W - Inches(1.6), Inches(0.8), caption, 18, color=GREY, align=PP_ALIGN.CENTER)
            add_text(s, Inches(0.8), Inches(6.4), W - Inches(1.6), Inches(0.5), ["Pearls Consulting"], 14, color=GREY, align=PP_ALIGN.CENTER)
            continue
        add_text(s, Inches(0.6), Inches(0.35), W - Inches(1.2), Inches(0.8), [title], 26, True)
        if hint == NO_SHOT:
            add_text(s, Inches(0.8), Inches(1.6), W - Inches(1.6), Inches(4.5), caption, 20)
            continue
        area = (Inches(0.6), Inches(1.15), W - Inches(1.2), Inches(4.45))
        files = shots.get(no) or []
        if not files:
            placeholder(s, *area, hint)
            missing.append(f"{no} ({hint}.png)")
        elif len(files) == 1:
            fit_picture(s, files[0], *area)
        else:  # side by side (e.g. EN | AR)
            gap = Inches(0.2)
            cw = int((area[2] - gap * (len(files) - 1)) / len(files))
            for i, f in enumerate(files):
                fit_picture(s, f, area[0] + i * (cw + gap), area[1], cw, area[3])
        add_text(s, Inches(0.6), Inches(5.7), W - Inches(1.2), Inches(1.7), caption, 13)
        s.notes_slide.notes_text_frame.text = f"Screenshot: {hint}.png"

    PRES_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"wrote {OUT}")
    print(f"screenshots found for slides: {sorted(shots) or 'none'}")
    if missing:
        print(f"placeholders left on slides: {missing}")


if __name__ == "__main__":
    build()
